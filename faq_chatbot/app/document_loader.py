import os
import re


def load_text(filepath):
    with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
        return f.read()


def load_pdf(filepath):
    import fitz
    doc = fitz.open(filepath)
    pages = []
    for page in doc:
        pages.append(page.get_text())
    doc.close()
    return '\n\n'.join(pages)


def load_docx(filepath):
    from docx import Document
    doc = Document(filepath)
    return '\n'.join(p.text for p in doc.paragraphs)


def load_pptx(filepath):
    from pptx import Presentation
    prs = Presentation(filepath)
    slides = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text)
        slides.append('\n'.join(texts))
    return '\n\n'.join(slides)


LOADERS = {
    '.txt': load_text,
    '.md': load_text,
    '.csv': load_text,
    '.json': load_text,
    '.xml': load_text,
    '.html': load_text,
    '.pdf': load_pdf,
    '.docx': load_docx,
    '.doc': load_docx,
    '.pptx': load_pptx,
    '.ppt': load_pptx,
}


def load_document(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    loader = LOADERS.get(ext)
    if not loader:
        raise ValueError(f"Unsupported file type: {ext}")
    text = loader(filepath)
    text = re.sub(r'\s+', ' ', text).strip()
    return text


def supported_extensions():
    return list(LOADERS.keys())
